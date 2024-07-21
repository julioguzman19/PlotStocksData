import os
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches

# Define the directory where your plots are saved
plot_dirs = ['Income_plots', 'Balance_plots', 'Cash_plots']

# Create a PowerPoint presentation object
prs = Presentation()

# Function to add a slide with plots
def add_slide_with_plots(prs, title, plot_files):
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    
    # Add images to the slide
    left = Inches(1)
    top = Inches(1.5)
    height = Inches(3.5)
    for plot_file in plot_files:
        slide.shapes.add_picture(plot_file, left, top, height=height)
        left += Inches(4.5)  # Adjust position for the next image
        if left > Inches(6.5):
            left = Inches(1)
            top += Inches(4)  # Move to the next row

# Function to extract the type of plot from the filename
def get_plot_type(filename):
    return '_'.join(filename.split('_')[1:])

# Collect all plots and group them by type
plot_files_by_type = defaultdict(list)
for plot_dir in plot_dirs:
    if os.path.exists(plot_dir):
        for filename in os.listdir(plot_dir):
            if filename.endswith('.png'):
                plot_type = get_plot_type(filename)
                plot_files_by_type[plot_type].append(os.path.join(plot_dir, filename))

# Create slides for each type of plot
for plot_type, plot_files in plot_files_by_type.items():
    add_slide_with_plots(prs, plot_type.replace('_', ' '), plot_files)

# Save the presentation
prs.save('Combined_Plots_Presentation.pptx')
